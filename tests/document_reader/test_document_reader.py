from __future__ import annotations

import builtins
import importlib.util
import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from contextlib import contextmanager
from pathlib import Path
from unittest import mock

from tools.document_reader.models import ReaderLimits, detect_file_type
from tools.document_reader.read_document import read_document

W_XMLNS = 'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'


@contextmanager
def _forced_import_error(*module_names: str):
    """Make `import <name>` raise ImportError for the given module names.

    Lets a test deterministically exercise a parser's standard-library
    fallback branch regardless of whether the optional dependency happens to
    be installed in the current environment.
    """
    real_import = builtins.__import__

    def fake_import(name, globals=None, locals=None, fromlist=(), level=0):  # noqa: A002
        if name in module_names:
            raise ImportError(f"forced ImportError for test: {name}")
        return real_import(name, globals, locals, fromlist, level)

    with mock.patch("builtins.__import__", side_effect=fake_import):
        yield


def _has_module(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


class DocumentReaderTests(unittest.TestCase):
    def test_file_type_detection(self) -> None:
        self.assertEqual(detect_file_type("example.txt"), "txt")
        self.assertEqual(detect_file_type("example.md"), "markdown")
        self.assertEqual(detect_file_type("example.xlsx"), "xlsx")
        with self.assertRaises(ValueError):
            detect_file_type("example.bin")

    def test_txt_parsing(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "note.txt"
            path.write_text("hello\nworld", encoding="utf-8")
            document = read_document(path)
            self.assertEqual(document.file_type, "txt")
            self.assertIn("hello", document.sections[0]["text"])

    def test_markdown_parsing(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "note.md"
            path.write_text("# Heading\n\nBody", encoding="utf-8")
            document = read_document(path)
            self.assertEqual(document.file_type, "markdown")
            self.assertIn("# Heading", document.sections[0]["text"])

    def test_json_parsing(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "data.json"
            path.write_text('{"a": 1}', encoding="utf-8")
            document = read_document(path)
            self.assertEqual(document.sections[0]["data"], {"a": 1})

    def test_json_with_utf8_bom_is_not_treated_as_malformed(self) -> None:
        # Regression test: a plain "utf-8" decode succeeds on BOM-prefixed
        # bytes without raising, but leaves the BOM character embedded in the
        # string, which then makes an otherwise well-formed JSON file fail to
        # parse with "Unexpected UTF-8 BOM". The BOM must be stripped before
        # json.loads ever sees the text.
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "data.json"
            path.write_bytes(b"\xef\xbb\xbf" + b'{"a": 1}')
            document = read_document(path)
            self.assertEqual(document.warnings, [])
            self.assertEqual(document.sections[0]["data"], {"a": 1})

    def test_large_json_omits_full_data(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "data.json"
            path.write_text(json.dumps({"items": ["abcdefghij"] * 50}), encoding="utf-8")
            document = read_document(path, ReaderLimits(max_chars=40))
            self.assertTrue(document.sections[0]["data_omitted"])
            self.assertNotIn("data", document.sections[0])
            self.assertIn("JSON preview truncated", " ".join(document.warnings))

    def test_malformed_json_keeps_truncation_warning(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "bad.json"
            path.write_text('{"a": "' + ("x" * 100), encoding="utf-8")
            document = read_document(path, ReaderLimits(max_chars=20))
            warnings = " ".join(document.warnings)
            self.assertIn("JSON raw text truncated", warnings)
            self.assertIn("Malformed JSON", warnings)

    def test_csv_parsing(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "data.csv"
            path.write_text("name,value\nA,1\nB,2\n", encoding="utf-8")
            document = read_document(path)
            self.assertEqual(document.tables[0]["headers"], ["name", "value"])
            self.assertEqual(len(document.tables[0]["rows"]), 2)

    def test_text_non_utf8_input_is_decoded_and_warned_not_corrupted(self) -> None:
        # Byte literals + \u escapes only (no embedded non-ASCII source
        # characters), so the test source is pure ASCII and proves correct
        # decoding regardless of how any editor, terminal, or diff tool
        # renders non-ASCII text. b"\xe9" and b"\xef" are the cp1252/latin-1
        # single-byte encodings of U+00E9 and U+00EF; "\u00e9" and "\u00ef"
        # are those same code points as explicit escapes on the Python side.
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "latin1.txt"
            path.write_bytes(b"caf\xe9 r\xe9sum\xe9 na\xefve")
            document = read_document(path)
            text = document.sections[0]["text"]
            self.assertEqual(text, "caf\u00e9 r\u00e9sum\u00e9 na\u00efve")
            self.assertNotIn("\ufffd", text)
            self.assertTrue(any("not valid UTF-8" in w for w in document.warnings))
            self.assertEqual(document.metadata["encoding"], "cp1252")

    def test_csv_non_utf8_input_is_decoded_and_warned_not_corrupted(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "latin1.csv"
            path.write_bytes(b"name,city\ncaf\xe9,na\xefve\n")
            document = read_document(path)
            self.assertEqual(document.tables[0]["rows"][0], ["caf\u00e9", "na\u00efve"])
            self.assertTrue(any("not valid UTF-8" in w for w in document.warnings))

    def test_xlsx_parsing_with_minimal_archive(self) -> None:
        # This fixture is deliberately a minimal, incomplete OOXML package
        # that only the hand-written fallback parser can read; force the
        # fallback branch so the test is not at the mercy of whether
        # openpyxl happens to be installed in the current environment.
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "book.xlsx"
            _write_minimal_xlsx(path)
            with _forced_import_error("openpyxl"):
                document = read_document(path, ReaderLimits(max_rows=10))
            self.assertEqual(document.file_type, "xlsx")
            self.assertEqual(document.sheets[0]["name"], "Revenue")
            self.assertEqual(document.sheets[0]["headers"], ["Month", "Amount"])
            self.assertEqual(document.sheets[0]["rows"][0][0], "Jan")
            self.assertEqual(str(document.sheets[0]["rows"][0][1]), "10")

    def test_pptx_fallback_respects_text_limit(self) -> None:
        # Same reasoning as the xlsx fallback test above: this fixture is an
        # incomplete OOXML package that only the fallback parser can read.
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "slides.pptx"
            _write_minimal_pptx(path, "A" * 80)
            with _forced_import_error("pptx"):
                document = read_document(path, ReaderLimits(max_chars=25))
            self.assertEqual(document.file_type, "pptx")
            self.assertEqual(len(document.slides[0]["text_blocks"][0]), 25)
            self.assertIn("PPTX text truncated", " ".join(document.warnings))

    def test_pptx_fallback_slide_count_is_capped(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "many_slides.pptx"
            _write_minimal_pptx_multi(path, slide_count=3)
            with _forced_import_error("pptx"):
                document = read_document(path, ReaderLimits(max_slides=2))
            self.assertEqual(len(document.slides), 2)
            self.assertTrue(any("PPTX slides truncated" in w for w in document.warnings))

    @unittest.skipUnless(_has_module("pptx"), "python-pptx not installed")
    def test_pptx_real_library_parsing(self) -> None:
        import pptx as pptx_lib

        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "real.pptx"
            presentation = pptx_lib.Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Hello"
            slide.placeholders[1].text = "World"
            presentation.save(path)

            document = read_document(path)
            self.assertEqual(document.file_type, "pptx")
            self.assertEqual(document.slides[0]["title"], "Hello")
            self.assertIn("World", document.slides[0]["text_blocks"])

    def test_docx_fallback_does_not_duplicate_nested_table_or_cell_paragraphs(self) -> None:
        # A table nested inside another table's cell previously got counted
        # twice (once via the parent cell's flattened text, once as its own
        # top-level table via a recursive ".//w:tbl" search), and cell
        # paragraphs leaked into the top-level section list via ".//w:p".
        xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<w:document {W_XMLNS}>
  <w:body>
    <w:p><w:r><w:t>Intro paragraph</w:t></w:r></w:p>
    <w:tbl>
      <w:tr>
        <w:tc>
          <w:p><w:r><w:t>Outer A1</w:t></w:r></w:p>
          <w:tbl>
            <w:tr><w:tc><w:p><w:r><w:t>Nested A1</w:t></w:r></w:p></w:tc></w:tr>
          </w:tbl>
        </w:tc>
      </w:tr>
    </w:tbl>
    <w:p><w:r><w:t>Outro paragraph</w:t></w:r></w:p>
  </w:body>
</w:document>"""
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "nested.docx"
            with zipfile.ZipFile(path, "w") as archive:
                archive.writestr("word/document.xml", xml)
            with _forced_import_error("docx"):
                document = read_document(path)
            self.assertEqual([s["text"] for s in document.sections], ["Intro paragraph", "Outro paragraph"])
            self.assertEqual(len(document.tables), 1)
            self.assertIn("Outer A1", document.tables[0]["headers"][0])
            self.assertIn("Nested A1", document.tables[0]["headers"][0])

    def test_docx_fallback_table_count_is_capped(self) -> None:
        tables_xml = "".join(
            f'<w:tbl><w:tr><w:tc><w:p><w:r><w:t>T{i}</w:t></w:r></w:p></w:tc></w:tr></w:tbl>' for i in range(5)
        )
        xml = f'<?xml version="1.0" encoding="UTF-8"?><w:document {W_XMLNS}><w:body>{tables_xml}</w:body></w:document>'
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "many_tables.docx"
            with zipfile.ZipFile(path, "w") as archive:
                archive.writestr("word/document.xml", xml)
            with _forced_import_error("docx"):
                document = read_document(path, ReaderLimits(max_tables=2))
            self.assertEqual(len(document.tables), 2)
            self.assertTrue(any("DOCX tables truncated" in w for w in document.warnings))

    @unittest.skipUnless(_has_module("docx"), "python-docx not installed")
    def test_docx_real_library_parsing_and_table_cap(self) -> None:
        import docx as docx_lib

        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "real.docx"
            doc = docx_lib.Document()
            doc.add_heading("Title", level=1)
            doc.add_paragraph("Body text")
            for i in range(3):
                doc.add_table(rows=1, cols=1).cell(0, 0).text = f"T{i}"
            doc.save(path)

            document = read_document(path, ReaderLimits(max_tables=2))
            self.assertEqual(document.file_type, "docx")
            self.assertTrue(any(s["type"] == "heading" for s in document.sections))
            self.assertEqual(len(document.tables), 2)
            self.assertTrue(any("DOCX tables truncated" in w for w in document.warnings))

    def test_pdf_no_parser_installed(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "doc.pdf"
            path.write_bytes(b"%PDF-1.4\n%%EOF")
            with _forced_import_error("docling.document_converter", "pymupdf"):
                document = read_document(path)
            self.assertEqual(document.file_type, "pdf")
            self.assertIn("No PDF parser is installed", " ".join(document.warnings))

    @unittest.skipUnless(_has_module("pymupdf"), "pymupdf not installed")
    def test_pdf_pymupdf_parsing(self) -> None:
        import pymupdf

        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "doc.pdf"
            pdf_doc = pymupdf.open()
            page = pdf_doc.new_page()
            page.insert_text((72, 72), "Hello PDF world")
            pdf_doc.save(path)
            pdf_doc.close()

            document = read_document(path)
            self.assertEqual(document.file_type, "pdf")
            self.assertEqual(document.metadata.get("parser"), "pymupdf")
            self.assertIn("Hello PDF world", document.sections[0]["text"])

    @unittest.skipUnless(_has_module("pymupdf"), "pymupdf not installed")
    def test_cli_pdf_stdout_is_pure_json(self) -> None:
        # Regression test: importing the deprecated `fitz` alias prints a
        # banner straight to stdout, which corrupts the JSON contract for any
        # caller piping this CLI's stdout into a JSON parser.
        import pymupdf

        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "doc.pdf"
            pdf_doc = pymupdf.open()
            pdf_doc.new_page()
            pdf_doc.save(path)
            pdf_doc.close()

            result = subprocess.run(
                [sys.executable, "tools/document_reader/read_document.py", str(path)],
                check=True,
                capture_output=True,
                text=True,
            )
            data = json.loads(result.stdout)
            self.assertEqual(data["file_type"], "pdf")

    def test_cli_exit_code_is_nonzero_for_missing_file(self) -> None:
        result = subprocess.run(
            [sys.executable, "tools/document_reader/read_document.py", "does/not/exist.pdf"],
            capture_output=True,
            text=True,
        )
        self.assertEqual(result.returncode, 2)

    def test_cli_exit_code_is_nonzero_for_unsupported_extension(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "file.xyz"
            path.write_text("data", encoding="utf-8")
            result = subprocess.run(
                [sys.executable, "tools/document_reader/read_document.py", str(path)],
                capture_output=True,
                text=True,
            )
            self.assertEqual(result.returncode, 2)

    def test_cli_exit_code_is_zero_when_parsed_with_warnings(self) -> None:
        # Truncation and other soft warnings are not failures: a Document was
        # still produced, so the exit code contract says 0.
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "note.txt"
            path.write_text("hello world", encoding="utf-8")
            result = subprocess.run(
                [sys.executable, "tools/document_reader/read_document.py", str(path), "--max-chars", "3"],
                capture_output=True,
                text=True,
            )
            self.assertEqual(result.returncode, 0)
            data = json.loads(result.stdout)
            self.assertTrue(data["warnings"])

    def test_malformed_file_handling(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "bad.xlsx"
            path.write_text("not a workbook", encoding="utf-8")
            document = read_document(path)
            self.assertTrue(document.warnings)
            self.assertIn("Could not parse workbook", document.warnings[0])

    def test_cli_json_output(self) -> None:
        with tempfile.TemporaryDirectory() as temp:
            path = Path(temp) / "note.txt"
            path.write_text("hello", encoding="utf-8")
            result = subprocess.run(
                [sys.executable, "tools/document_reader/read_document.py", str(path)],
                check=True,
                capture_output=True,
                text=True,
            )
            data = json.loads(result.stdout)
            self.assertEqual(data["file_type"], "txt")


def _write_minimal_xlsx(path: Path) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"></Types>""")
        archive.writestr(
            "xl/workbook.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets><sheet name="Revenue" sheetId="1" r:id="rId1"/></sheets>
</workbook>""",
        )
        archive.writestr(
            "xl/_rels/workbook.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "xl/sharedStrings.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <si><t>Month</t></si><si><t>Amount</t></si><si><t>Jan</t></si>
</sst>""",
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <dimension ref="A1:B2"/>
  <sheetData>
    <row r="1"><c r="A1" t="s"><v>0</v></c><c r="B1" t="s"><v>1</v></c></row>
    <row r="2"><c r="A2" t="s"><v>2</v></c><c r="B2"><v>10</v></c></row>
  </sheetData>
</worksheet>""",
        )


def _write_minimal_pptx(path: Path, text: str) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"></Types>""")
        archive.writestr(
            "ppt/slides/slide1.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld><p:spTree><p:sp><p:txBody><a:p><a:r><a:t>{text}</a:t></a:r></a:p></p:txBody></p:sp></p:spTree></p:cSld>
</p:sld>""",
        )


def _write_minimal_pptx_multi(path: Path, slide_count: int) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr("[Content_Types].xml", """<?xml version="1.0" encoding="UTF-8"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"></Types>""")
        for i in range(1, slide_count + 1):
            archive.writestr(
                f"ppt/slides/slide{i}.xml",
                f"""<?xml version="1.0" encoding="UTF-8"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:cSld><p:spTree><p:sp><p:txBody><a:p><a:r><a:t>Slide {i}</a:t></a:r></a:p></p:txBody></p:sp></p:spTree></p:cSld>
</p:sld>""",
            )


if __name__ == "__main__":
    unittest.main()
