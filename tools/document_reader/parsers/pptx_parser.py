from __future__ import annotations

import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from tools.document_reader.models import Document, ReaderLimits

A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def parse_pptx(path: Path, limits: ReaderLimits) -> Document:
    try:
        import pptx  # type: ignore
    except ImportError:
        return _parse_pptx_zip(path, limits)

    warnings: list[str] = []
    try:
        presentation = pptx.Presentation(path)
    except Exception as exc:
        return Document(path.name, "pptx", warnings=[f"Could not parse PPTX: {exc}"])
    slides = []
    chars = 0
    total_slides = len(presentation.slides)
    for index, slide in enumerate(presentation.slides, start=1):
        if index > limits.max_slides:
            warnings.append(f"PPTX slides truncated to {limits.max_slides} from {total_slides}.")
            break
        text_blocks = []
        tables = []
        title = slide.shapes.title.text if slide.shapes.title else ""
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    limited_text, chars, truncated = _append_limited_text(text, chars, limits.max_chars)
                    if limited_text:
                        text_blocks.append(limited_text)
                    if truncated:
                        warnings.append(f"PPTX text truncated to {limits.max_chars} characters.")
                        break
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows[: limits.table_preview_rows]]
                tables.append({"headers": rows[0] if rows else [], "rows": rows[1:]})
        slides.append({"slide_number": index, "title": title, "text_blocks": text_blocks, "tables": tables})
        if chars >= limits.max_chars:
            break
    if not slides:
        warnings.append("PPTX appears empty.")
    return Document(path.name, "pptx", slides=slides, warnings=warnings)


def _parse_pptx_zip(path: Path, limits: ReaderLimits) -> Document:
    warnings = ["python-pptx is not installed; using limited PPTX XML parser."]
    try:
        with zipfile.ZipFile(path) as archive:
            slide_paths = sorted(
                name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            )
            slides = []
            chars = 0
            total_slides = len(slide_paths)
            for index, slide_path in enumerate(slide_paths, start=1):
                if index > limits.max_slides:
                    warnings.append(f"PPTX slides truncated to {limits.max_slides} from {total_slides}.")
                    break
                root = ET.fromstring(archive.read(slide_path))
                texts = []
                for node in root.findall(".//a:t", A_NS):
                    if not node.text:
                        continue
                    limited_text, chars, truncated = _append_limited_text(node.text, chars, limits.max_chars)
                    if limited_text:
                        texts.append(limited_text)
                    if truncated:
                        warnings.append(f"PPTX text truncated to {limits.max_chars} characters.")
                        break
                slides.append({"slide_number": index, "title": texts[0] if texts else "", "text_blocks": texts, "tables": []})
                if chars >= limits.max_chars:
                    break
    except zipfile.BadZipFile:
        return Document(path.name, "pptx", warnings=["Could not parse PPTX: invalid or corrupt archive."])
    except Exception as exc:
        return Document(path.name, "pptx", warnings=[f"Could not parse PPTX: {exc}"])
    if not slides:
        warnings.append("PPTX appears empty.")
    return Document(path.name, "pptx", slides=slides, warnings=warnings)


def _append_limited_text(text: str, chars: int, max_chars: int) -> tuple[str, int, bool]:
    remaining = max_chars - chars
    if remaining <= 0:
        return "", chars, True
    if len(text) > remaining:
        return text[:remaining], max_chars, True
    return text, chars + len(text), False
